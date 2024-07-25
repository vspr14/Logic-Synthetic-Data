from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
prs = Presentation()

# Function to add a slide with a title and content
def add_slide(prs, title_text, content_text, layout_index=1):
    slide_layout = prs.slide_layouts[layout_index]  # Title and Content Layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)

    content.text = content_text
    content.text_frame.paragraphs[0].font.size = Pt(24)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 153, 51)

# Slide 1: Problem Statement and Goal
slide_layout = prs.slide_layouts[0]  # Title Slide Layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Problem Statement and Goal"
title.text_frame.paragraphs[0].font.size = Pt(44)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)

subtitle.text = (
    "Problem Statement:\n"
    "Retail businesses often struggle with generating sufficient data for training machine learning models, "
    "especially for detecting valid and invalid transactions.\n\n"
    "Goal:\n"
    "Develop a tool that allows users to upload an Excel file and generate synthetic data for retail applications. "
    "This synthetic data can be used to train ML models, improving their accuracy in detecting valid and invalid transactions."
)
subtitle.text_frame.paragraphs[0].font.size = Pt(24)
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 153, 51)

# Slide 2: Approach - Part 1
add_slide(prs, "Approach - Part 1", 
    "1. Data Upload and Preprocessing:\n"
    "- Users upload an Excel file containing retail transaction data.\n"
    "- The data is preprocessed to remove any missing or invalid entries.\n\n"
    "2. Data Filtering:\n"
    "- Users can filter the data based on store number and other criteria.\n"
    "- The filtered data is used for further processing."
)

# Slide 3: Approach - Part 2
add_slide(prs, "Approach - Part 2", 
    "3. Synthetic Data Generation:\n"
    "- Users can choose to add noise to the data to simulate real-world scenarios.\n"
    "- The tool uses Gaussian Copula Synthesizer to generate synthetic data based on the filtered data.\n\n"
    "4. Data Visualization and Metrics:\n"
    "- The tool provides visualizations to compare the original and synthetic data.\n"
    "- Metrics are calculated to evaluate the quality of the synthetic data."
)

# Slide 4: Benefits
add_slide(prs, "Benefits", 
    "1. Improved ML Model Accuracy:\n"
    "- Synthetic data helps in training machine learning models, improving their accuracy in detecting valid and invalid transactions.\n\n"
    "2. Data Privacy:\n"
    "- Synthetic data generation ensures that sensitive customer information is not exposed.\n\n"
    "3. Scalability:\n"
    "- The tool can generate large amounts of data, making it suitable for various retail applications.\n\n"
    "4. Flexibility:\n"
    "- Users can customize the data generation process by adding noise and setting constraints."
)

# Slide 5: Tech Stack
add_slide(prs, "Tech Stack", 
    "1. Python:\n"
    "- Core programming language used for development.\n\n"
    "2. Streamlit:\n"
    "- Used for building the web interface for data upload and visualization.\n\n"
    "3. SDV (Synthetic Data Vault):\n"
    "- Used for generating synthetic data using Gaussian Copula Synthesizer.\n\n"
    "4. Pandas:\n"
    "- Used for data manipulation and preprocessing.\n\n"
    "5. Plotly:\n"
    "- Used for creating interactive visualizations.\n\n"
    "6. Scikit-learn:\n"
    "- Used for implementing machine learning algorithms and metrics."
)

# Save the presentation
prs.save('Styled_Synthetic_Data_Generation_Presentation.pptx')
print("Styled presentation created successfully!")